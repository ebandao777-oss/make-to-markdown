#!/usr/bin/env python3
"""
Make-to-Markdown 智能转换引擎 (v3.6)

自动检测文件格式 → 检查依赖 → 缺则安装 → markitdown 转换 → 失败降级 → 后置清洗
确保：清晰的标题层级 + 纯净的文本内容 + 文档级摘要 + 结构化数据 + 完整的上下文

v2.8: 启动时检测 OS/版本/能力，按平台选择最佳执行路径 + 共享 platform_detect 模块
v2.7: uv 可用性显式检查 + win32com 错误提示改进
v2.5: markitdown extras 自动安装 + COM 异常安全加固 + .ppt/.doc 全链路容错
v2.3: 新增 .doc/.ppt 老格式自动预处理（Word/PowerPoint COM / LibreOffice → .docx/.pptx）

用法:
    python convert.py input.docx -o output.md
    python convert.py input.doc -o output.md
    python convert.py input.pdf --summary --breadcrumbs
    python convert.py ./docs/ -o ./output/  --batch
"""

import os
import sys
import re
import subprocess
import importlib.util
import argparse
import tempfile
import shutil
from pathlib import Path
from typing import Optional, Tuple, List, Dict

VERSION = "3.6"

# ── 启动时平台检测 ──
try:
    from platform_detect import detect as _platform_detect, best_office_tool
    _PINFO = _platform_detect()
except ImportError:
    _PINFO = None
    def best_office_tool(_pinfo=None):
        return "com" if sys.platform == "win32" else "libreoffice"


# ═══════════════════════════════════════════════════════════════
# 1. 格式→依赖映射
# ═══════════════════════════════════════════════════════════════

# (扩展名, pip 包名, import 检测模块, markitdown extra 名)
DEPENDENCY_MAP: Dict[str, Tuple[str, str, str]] = {
    '.docx': ('python-docx', 'docx', 'docx'),
    '.doc':  ('python-docx', 'docx', 'docx'),
    '.xlsx': ('openpyxl', 'openpyxl', 'xlsx'),
    '.xls':  ('xlrd', 'xlrd', 'xls'),
    '.pptx': ('python-pptx', 'pptx', 'pptx'),
    '.ppt':  ('python-pptx', 'pptx', 'pptx'),
    '.pdf':  ('pypdf', 'pypdf', 'pdf'),
    '.epub': ('ebooklib', 'ebooklib', 'epub'),
}

SUPPORTED_EXTS = set(DEPENDENCY_MAP.keys()) | {
    '.html', '.htm', '.csv', '.json', '.xml',
    '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp',
    '.mp3', '.wav', '.m4a', '.ogg', '.flac', '.zip',
}


def check_dependency(import_name: str) -> bool:
    """检查 Python 模块是否已安装"""
    return importlib.util.find_spec(import_name) is not None


def install_dependency(pkg_name: str) -> bool:
    """使用当前 Python 环境安装包"""
    try:
        print(f"  [INFO] 将执行: pip install {pkg_name}（从 PyPI 下载，约数 MB，不会上传本机数据）")
        print(f"  [依赖] 安装 {pkg_name} ...", end='', flush=True)
        result = subprocess.run(
            [sys.executable, '-m', 'pip', 'install', pkg_name, '-q'],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode == 0:
            print(' OK')
            return True
        print(f' 失败 ({result.stderr.strip()[-100:]})')
        return False
    except Exception as e:
        print(f' 失败 ({e})')
        return False


def ensure_dependencies(ext: str) -> bool:
    """确保指定文件格式的依赖已安装"""
    ext = ext.lower()
    if ext not in DEPENDENCY_MAP:
        return True  # 内置支持格式，无需额外依赖

    pkg_name, import_name, _ = DEPENDENCY_MAP[ext]
    if check_dependency(import_name):
        return True
    print(f"[依赖检测] 缺少 {import_name}，自动安装...")
    return install_dependency(pkg_name)


def ensure_markitdown_extras() -> bool:
    """确保 markitdown（uv tool 侧）安装了所有可选依赖。"""
    uv_bin = shutil.which('uv')
    if not uv_bin:
        print("  [uv] 未在 PATH 中找到 uv，跳过 markitdown extras 安装（不影响转换，依赖缺失时走降级）")
        return False

    # 检查 uv tool list 是否有 markitdown
    try:
        check = subprocess.run(
            [uv_bin, 'tool', 'list'], capture_output=True, text=True, timeout=10,
            encoding='utf-8', errors='replace',
        )
        if 'markitdown' not in (check.stdout or '') or check.returncode != 0:
            _install_markitdown_with_extras(uv_bin)
    except Exception:
        _install_markitdown_with_extras(uv_bin)

    return True


def _install_markitdown_with_extras(uv_bin: str) -> None:
    """用 uv tool install 安装 markitdown 及全部可选依赖"""
    extra_pkgs = [
        'lxml', 'mammoth',          # [docx]
        'python-pptx',              # [pptx]
        'openpyxl', 'pandas',       # [xlsx]
        'pandas', 'xlrd',           # [xls]
        'pdfminer-six', 'pdfplumber',  # [pdf]
        'olefile',                  # [outlook]
        'ebooklib',                 # [epub]
        'youtube-transcript-api',   # [youtube-transcription]
        'pydub',                    # [audio-transcription]
        'speechrecognition',        # [audio-transcription]
    ]
    # 去重并安装
    unique_pkgs = list(dict.fromkeys(extra_pkgs))
    try:
        subprocess.run(
            [uv_bin, 'tool', 'install', 'markitdown', '--force']
            + [f'--with={pkg}' for pkg in unique_pkgs],
            capture_output=True, text=True, timeout=180,
            encoding='utf-8', errors='replace',
        )
        print("  [markitdown extras] 已安装全部可选依赖")
    except Exception as e:
        print(f"  [markitdown extras] 安装失败: {e}")


# ═══════════════════════════════════════════════════════════════
# 2. markitdown 转换
# ═══════════════════════════════════════════════════════════════

def find_markitdown() -> Optional[str]:
    """查找 markitdown 可执行文件（~/.local/bin 或 uv tool dir）"""
    if sys.platform != 'win32':
        candidates = [
            os.path.join(os.path.expanduser('~'), '.local', 'bin', 'markitdown'),
        ]
    else:
        candidates = [
            os.path.join(os.path.expanduser('~'), '.local', 'bin', 'markitdown.exe'),
            os.path.join(os.path.expanduser('~'), '.local', 'bin', 'markitdown'),
        ]
    # 追加 uv tool dir
    uv_bin = shutil.which('uv')
    if uv_bin:
        try:
            result = subprocess.run(
                [uv_bin, 'tool', 'dir'], capture_output=True, text=True, timeout=10,
                encoding='utf-8', errors='replace',
            )
            tool_dir = (result.stdout or '').strip()
            if tool_dir:
                if sys.platform == 'win32':
                    candidates.append(os.path.join(tool_dir, 'markitdown.exe'))
                candidates.append(os.path.join(tool_dir, 'markitdown'))
        except Exception:
            pass
    for c in candidates:
        if os.path.exists(c):
            return c
    return None


def convert_with_markitdown(input_path: str, output_path: str, extra_pkg: Optional[str] = None) -> bool:
    """用 markitdown CLI 转换文件。extra_pkg 为 uvx 模式传递 --with 参数。"""
    bin_path = find_markitdown()
    if not bin_path:
        cmd = ['uvx']
        if extra_pkg:
            cmd.extend(['--with', extra_pkg])
        cmd.extend(['markitdown', input_path, '-o', output_path])
    else:
        cmd = [bin_path, input_path, '-o', output_path]

    try:
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
        if result.returncode == 0 and os.path.exists(output_path):
            return True
        # 检查是否为依赖缺失错误
        stderr = result.stderr or ''
        if 'MissingDependencyException' in stderr or 'dependencies needed' in stderr:
            print(f"  [markitdown] 依赖缺失: {stderr.strip()[-200:]}")
        else:
            print(f"  [markitdown] 转换失败: {stderr.strip()[-200:]}")
        return False
    except subprocess.TimeoutExpired:
        print(f"  [markitdown] 超时")
        return False
    except Exception as e:
        print(f"  [markitdown] 异常: {e}")
        return False


# ═══════════════════════════════════════════════════════════════
# 3. 降级转换器（markitdown 失败时兜底）
# ═══════════════════════════════════════════════════════════════

def convert_docx_fallback(input_path: str) -> str:
    """python-docx 原生转换 Word → Markdown"""
    try:
        from docx import Document
    except ImportError:
        return ''

    doc = Document(input_path)
    md_lines = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            md_lines.append('')
            continue

        style_name = para.style.name if para.style else ''
        fs = para.runs[0].font.size if para.runs else None
        fs_pt = fs / 12700 if fs else 0

        # 标题识别：样式 + 字号启发式
        if style_name == 'Title':
            md_lines.append(f'# {text}')
        elif style_name == 'Heading 1' or style_name.startswith('Heading 1'):
            md_lines.append(f'# {text}')
        elif style_name == 'Heading 2' or (style_name == 'List Paragraph' and fs_pt >= 21):
            md_lines.append(f'## {text}')
        elif style_name.startswith('Heading 3') or (style_name == 'List Paragraph' and 16 <= fs_pt < 21):
            md_lines.append(f'### {text}')
        elif style_name.startswith('Heading'):
            level = int(style_name.split()[-1]) if style_name.split()[-1].isdigit() else 3
            md_lines.append(f'{"#" * level} {text}')
        else:
            # 粗体短文本 → 可能是内嵌标题
            runs = para.runs
            if runs and len(text) < 50 and all(r.bold for r in runs if r.text.strip()):
                if re.match(r'^\d+[\.\、\s)]', text):
                    md_lines.append(f'### {text}')
                else:
                    md_lines.append(f'**{text}**')
            else:
                md_lines.append(text)

    # 表格
    for table in doc.tables:
        md_lines.append('')
        rows = [[c.text.strip().replace('\n', ' ') for c in row.cells] for row in table.rows]
        if len(rows) >= 2:
            md_lines.append('| ' + ' | '.join(rows[0]) + ' |')
            md_lines.append('| ' + ' | '.join('---' for _ in rows[0]) + ' |')
            for row in rows[1:]:
                md_lines.append('| ' + ' | '.join(row) + ' |')

    return '\n'.join(md_lines)


def convert_xlsx_fallback(input_path: str) -> str:
    """openpyxl 原生转换 Excel → Markdown 表格"""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ''

    wb = load_workbook(input_path, data_only=True)
    md_lines = [f'# {Path(input_path).stem}\n']

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        md_lines.append(f'## {sheet_name}\n')

        # 获取有效数据范围
        rows = list(ws.iter_rows(values_only=True))
        rows = [r for r in rows if any(c is not None for c in r)]

        if not rows:
            continue

        max_cols = max(len(r) for r in rows)
        # 补全列
        padded = [list(r) + [''] * (max_cols - len(r)) for r in rows]

        header = '| ' + ' | '.join(str(c or '') for c in padded[0]) + ' |'
        sep = '| ' + ' | '.join('---' for _ in range(max_cols)) + ' |'
        md_lines.append(header)
        md_lines.append(sep)

        for row in padded[1:]:
            md_lines.append('| ' + ' | '.join(str(c or '') for c in row) + ' |')
        md_lines.append('')

    return '\n'.join(md_lines)


def convert_pptx_fallback(input_path: str) -> str:
    """python-pptx 原生转换 PPT → Markdown"""
    try:
        from pptx import Presentation
    except ImportError:
        return ''

    prs = Presentation(input_path)
    md_lines = [f'# {Path(input_path).stem}\n']

    for i, slide in enumerate(prs.slides, 1):
        md_lines.append(f'## 幻灯片 {i}\n')
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        level = para.level if para.level else 0
                        prefix = '  ' * level + '- ' if level > 0 else ''
                        md_lines.append(f'{prefix}{text}')
            if shape.has_table:
                table = shape.table
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                if rows:
                    md_lines.append('| ' + ' | '.join(rows[0]) + ' |')
                    md_lines.append('| ' + ' | '.join('---' for _ in rows[0]) + ' |')
                    for row in rows[1:]:
                        md_lines.append('| ' + ' | '.join(row) + ' |')
        md_lines.append('')

    return '\n'.join(md_lines)


# ═══════════════════════════════════════════════════════════════
# 2.5 旧格式预处理 .doc/.ppt → .docx/.pptx
# ═══════════════════════════════════════════════════════════════

# 旧格式 → 新格式映射
LEGACY_TO_MODERN = {
    '.doc': '.docx',
    '.ppt': '.pptx',
}


def _legacy_via_win32com(input_path: str, output_dir: str, target_ext: str) -> Optional[str]:
    """通过 Office COM 自动化将旧格式另存为新格式"""
    if sys.platform != 'win32':
        return None
    ext = Path(input_path).suffix.lower()
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        print(f"  [预处理] pywin32 未安装，无法使用 Office COM。安装: pip install pywin32")
        return None

    pythoncom.CoInitialize()
    app = None
    output_path = None
    try:
        if ext == '.doc':
            app = win32com.client.Dispatch("Word.Application")
            for attr, val in [("Visible", False), ("DisplayAlerts", 0)]:
                try:
                    setattr(app, attr, val)
                except Exception:
                    pass
            doc = app.Documents.Open(os.path.abspath(input_path), ReadOnly=True)
            output_path = os.path.join(output_dir, Path(input_path).stem + target_ext)
            doc.SaveAs2(os.path.abspath(output_path), FileFormat=16)  # 16=wdFormatXMLDocument
            doc.Close()
        elif ext == '.ppt':
            app = win32com.client.Dispatch("PowerPoint.Application")
            for attr, val in [("Visible", False), ("DisplayAlerts", 0)]:
                try:
                    setattr(app, attr, val)
                except Exception:
                    pass
            pres = app.Presentations.Open(os.path.abspath(input_path), ReadOnly=True)
            output_path = os.path.join(output_dir, Path(input_path).stem + target_ext)
            pres.SaveAs(os.path.abspath(output_path), FileFormat=24)  # 24=ppSaveAsOpenXMLPresentation
            pres.Close()
        else:
            return None
    except Exception as e:
        app_label = 'Word' if ext == '.doc' else 'PowerPoint'
        print(f"  [预处理] {app_label} COM 不可用 ({e})")
        return None
    finally:
        try:
            if app:
                app.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()

    if output_path and os.path.exists(output_path) and os.path.getsize(output_path) > 0:
        app_label = 'Word' if ext == '.doc' else 'PowerPoint'
        print(f"  [预处理] {app_label} COM: {ext} → {target_ext} OK")
        return output_path
    return None


def _legacy_via_libreoffice(input_path: str, output_dir: str, target_ext: str) -> Optional[str]:
    """通过 LibreOffice 命令行将旧格式转为新格式。
    优先级: PATH 中的 soffice → 环境变量 $SOFFICE_PATH → 常见安装路径 → 'soffice' 兜底"""
    ext = Path(input_path).suffix.lower()
    convert_to = target_ext.lstrip('.')
    soffice_paths = []

    # 1. PATH 中查找（跨平台首选）
    soffice_which = shutil.which("soffice")
    if soffice_which:
        soffice_paths.append(soffice_which)

    # 2. 环境变量 $SOFFICE_PATH（用户自定义路径）
    env_path = os.environ.get("SOFFICE_PATH")
    if env_path and os.path.exists(env_path) and env_path not in soffice_paths:
        soffice_paths.append(env_path)

    # 3. 常见安装路径（兜底）
    common_paths = []
    if sys.platform == 'win32':
        common_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
    else:
        common_paths = [
            "/usr/bin/soffice",
            "/usr/lib/libreoffice/program/soffice",
            "/usr/lib64/libreoffice/program/soffice",
            "/opt/libreoffice/program/soffice",
            "/snap/bin/libreoffice.soffice",
        ]
    for p in common_paths:
        if os.path.exists(p) and p not in soffice_paths:
            soffice_paths.append(p)

    # 4. 裸名兜底
    if "soffice" not in soffice_paths:
        soffice_paths.append("soffice")

    for soffice in soffice_paths:
        if shutil.which(soffice) or os.path.exists(soffice):
            try:
                abs_input = os.path.abspath(input_path)
                abs_output = os.path.abspath(output_dir)
                result = subprocess.run(
                    [soffice, '--headless', '--convert-to', convert_to,
                     '--outdir', abs_output, abs_input],
                    capture_output=True, text=True, timeout=60,
                )
                output_path = os.path.join(output_dir, Path(input_path).stem + target_ext)
                if result.returncode == 0 and os.path.exists(output_path):
                    print(f"  [预处理] LibreOffice: {ext} → {target_ext} OK")
                    return output_path
            except Exception as e:
                print(f"  [预处理] LibreOffice ({soffice}) 失败 ({e})")
                continue
    return None


def preprocess_legacy(input_path: str) -> Optional[str]:
    """
    将 .doc / .ppt 老格式转换为 .docx / .pptx。
    优先级（按平台能力选择）: Office COM → LibreOffice CLI。
    返回新文件路径，失败返回 None。
    """
    ext = Path(input_path).suffix.lower()
    target_ext = LEGACY_TO_MODERN.get(ext)
    if not target_ext:
        return None

    tmpdir = tempfile.mkdtemp(prefix='mtm_convert_')

    # 按平台能力选择最佳转换工具
    tool = best_office_tool(_PINFO)
    if tool == "com":
        result = _legacy_via_win32com(input_path, tmpdir, target_ext)
        if not result:
            result = _legacy_via_libreoffice(input_path, tmpdir, target_ext)
    elif tool == "libreoffice":
        result = _legacy_via_libreoffice(input_path, tmpdir, target_ext)
    else:
        result = None
    if not result:
        try:
            shutil.rmtree(tmpdir, ignore_errors=True)
        except Exception:
            pass
    return result


FALLBACK_MAP = {
    '.docx': convert_docx_fallback,
    '.doc': convert_docx_fallback,
    '.xlsx': convert_xlsx_fallback,
    '.xls': convert_xlsx_fallback,
    '.pptx': convert_pptx_fallback,
    '.ppt': convert_pptx_fallback,
}


# ═══════════════════════════════════════════════════════════════
# 4. 后置清洗（内联轻量版，避免 import 开销）
# ═══════════════════════════════════════════════════════════════

CLEANUP_PATTERNS = [
    (re.compile(r'^\s*\[Generated by[^\]]*\]\s*$', re.MULTILINE), ''),
    (re.compile(r'^\s*\*Converted by markitdown[^*]*\*\s*$', re.MULTILINE), ''),
    (re.compile(r'^\s*Page\s+\d+\s+of\s+\d+\s*$', re.MULTILINE | re.IGNORECASE), ''),
    (re.compile(r'^\s*\[PAGE\s*\d+\]\s*$', re.MULTILINE | re.IGNORECASE), ''),
    (re.compile(r'^\s*Confidential[\s\-–—].*$', re.MULTILINE | re.IGNORECASE), ''),
    (re.compile(r'^\s*DRAFT\s*.*$', re.MULTILINE | re.IGNORECASE), ''),
    (re.compile(r'^\s*Internal Use Only\s*$', re.MULTILINE | re.IGNORECASE), ''),
    (re.compile(r'^\s*For Review Only\s*$', re.MULTILINE | re.IGNORECASE), ''),
    (re.compile(r'^\s*\[Created with[^\]]*\]\s*$', re.MULTILINE), ''),
    (re.compile(r'^\s*Powered by[^\n]*markitdown[^\n]*$', re.MULTILINE), ''),
    (re.compile(r'^\s*---\s*Page\s+\d+\s*---\s*$', re.MULTILINE | re.IGNORECASE), ''),
    (re.compile(r'^\s*\u00a9\s*\d{4}.*$', re.MULTILINE), ''),
    (re.compile(r'^\s*All Rights Reserved\.?\s*$', re.MULTILINE), ''),
    (re.compile(r'\n{3,}'), '\n\n'),
    (re.compile(r'[ \t]+$', re.MULTILINE), ''),
]


def fix_heading_hierarchy(content: str) -> str:
    """修复标题层级：多 H1 降级 + 越级插入过渡标题"""
    lines = content.split('\n')
    headings = [(i, len(m.group(1)), m.group(2).strip())
                for i, line in enumerate(lines)
                if (m := re.match(r'^(#{1,6})\s+(.+)', line))]
    if not headings:
        return content

    actions = []  # (line_idx, 'demote' | 'insert', ...)
    first_h1 = True
    for idx, level, text in headings:
        if level == 1:
            if not first_h1:
                actions.append((idx, 'demote', text))
            first_h1 = False

    prev_level = 1
    for idx, level, text in headings:
        jump = level - prev_level
        if jump > 1:
            actions.append((idx, 'insert_before', level - 1, text))
        prev_level = level

    actions.sort(key=lambda a: a[0], reverse=True)
    for act in actions:
        if act[1] == 'demote':
            idx, _, text = act
            lines[idx] = f"## {text}"
        elif act[1] == 'insert_before':
            idx, _, target_lvl, ref = act
            prefix = '#' * target_lvl
            lines.insert(idx, f"\n{prefix} {ref}（概述）\n")

    return '\n'.join(lines)


def fix_tables(content: str) -> str:
    """为缺失分隔符的表格补全 |---|---|---| 行"""
    lines = content.split('\n')
    i = 0
    in_table = False
    while i < len(lines) - 1:
        current = lines[i].strip()
        if current.startswith('|') and current.endswith('|'):
            cols = [c.strip() for c in current.split('|')[1:-1]]
            if len(cols) >= 2:
                next_line = lines[i + 1].strip() if i + 1 < len(lines) else ''
                if re.match(r'^\|[\s\-:]+\|', next_line):
                    in_table = True
                    i += 2
                    continue
                if not next_line.startswith('|') or next_line.startswith('>'):
                    in_table = False
                    i += 1
                    continue
                if not in_table:
                    sep = '| ' + ' | '.join('---' for _ in cols) + ' |'
                    lines.insert(i + 1, sep)
                    in_table = True
                    i += 2
                    continue
        else:
            in_table = False
        i += 1
    return '\n'.join(lines)


def post_clean(content: str, with_summary: bool = True) -> str:
    """内联后置清洗"""
    for pat, rep in CLEANUP_PATTERNS:
        content = pat.sub(rep, content)

    content = fix_heading_hierarchy(content)
    content = fix_tables(content)

    if with_summary:
        headings = re.findall(r'^(#{1,3})\s+(.+)$', content, re.MULTILINE)
        h1 = [t for l, t in headings if l == '#']
        h2 = [t for l, t in headings if l == '##']
        tables = len(re.findall(r'^\| ---', content, re.MULTILINE))
        title = h1[0] if h1 else (h2[0] if h2 else Path('untitled').stem)

        summary = f"本文档「{title}」，涵盖 {len(h2)} 个章节"
        if tables:
            summary += f"，包含 {tables} 个数据表格"
        summary += "。"
        content = f"> **文档摘要**\n> {summary}\n\n---\n\n{content}"

    return content.strip() + '\n'


# ═══════════════════════════════════════════════════════════════
# 5. 主流程
# ═══════════════════════════════════════════════════════════════

def convert_file(input_path: str, output_path: str, with_summary: bool = True) -> Tuple[bool, str]:
    """
    智能转换单个文件。
    返回 (成功, 消息)
    """
    original_ext = Path(input_path).suffix.lower()
    ext = original_ext
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    docx_temp = None  # .doc 预处理产生的临时 .docx

    # Step 0: 旧格式预处理 .doc/.ppt → .docx/.pptx
    if original_ext in LEGACY_TO_MODERN:
        target_ext = LEGACY_TO_MODERN[original_ext]
        print(f"  [检测] {original_ext} 旧格式，尝试预处理为 {target_ext} ...")
        docx_temp = preprocess_legacy(input_path)
        if docx_temp:
            input_path = docx_temp
            ext = target_ext
            print(f"  [预处理] 使用中间文件: {input_path}")
        else:
            return False, (f"{original_ext} 旧格式无法处理。请安装 Microsoft Office "
                           f"或 LibreOffice 以启用自动转换，"
                           f"或手动将文件另存为 {target_ext} 格式。")

    # Step 1: 依赖检测 & 自动安装
    if not ensure_dependencies(ext):
        return False, f"无法安装 {ext} 所需依赖"

    # Step 1.5: 确保 markitdown extras（避免 MissingDependencyException）
    if ext in DEPENDENCY_MAP:
        ensure_markitdown_extras()

    # Step 2: 尝试 markitdown
    extra_pkg = DEPENDENCY_MAP[ext][0] if ext in DEPENDENCY_MAP else None
    markitdown_ok = convert_with_markitdown(input_path, output_path, extra_pkg=extra_pkg)

    content = ''
    if markitdown_ok:
        with open(output_path, 'r', encoding='utf-8') as f:
            content = f.read()
    elif ext in FALLBACK_MAP:
        # Step 3: 降级转换
        print(f"  [降级] markitdown 失败，使用原生转换器...")
        content = FALLBACK_MAP[ext](input_path)
        if content:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(content)
        else:
            return False, f"{ext} 格式转换失败（markitdown 和原生转换器均不可用）"
    else:
        return False, f"{ext} 格式转换失败且无降级方案"

    # Step 4: 后置清洗
    content = post_clean(content, with_summary=with_summary)
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(content)

    # 统计
    headings = re.findall(r'^(#{1,3})\s+(.+)$', content, re.MULTILINE)
    h1_c, h2_c, h3_c = 0, 0, 0
    for l, _ in headings:
        if l == '#': h1_c += 1
        elif l == '##': h2_c += 1
        elif l == '###': h3_c += 1
    tables = len(re.findall(r'^\| ---', content, re.MULTILINE))
    method = 'markitdown' if markitdown_ok else '原生降级'

    msg = (f"转换完成 [{method}] | H1={h1_c} H2={h2_c} H3={h3_c} | 表格={tables} | "
           f"{os.path.getsize(output_path)} bytes")

    # 清理 .doc 预处理产生的临时 .docx
    if docx_temp:
        try:
            tmpdir = os.path.dirname(docx_temp)
            shutil.rmtree(tmpdir, ignore_errors=True)
        except Exception:
            pass

    return True, msg


def main():
    parser = argparse.ArgumentParser(
        description='Make-to-Markdown 智能转换引擎',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog='''
示例:
  python convert.py report.docx -o output.md
  python convert.py old_doc.doc -o output.md
  python convert.py data.xlsx --summary
  python convert.py slide.pptx -o slide.md
        ''',
    )
    parser.add_argument('input', help='输入文件路径')
    parser.add_argument('--output', '-o', default=None, help='输出文件路径')
    parser.add_argument('--no-summary', action='store_true', help='跳过文档摘要')
    parser.add_argument('--quiet', '-q', action='store_true', help='静默模式')
    parser.add_argument('--version', '-V', action='version',
                        version=f'Make-to-Markdown convert.py v{VERSION}')
    args = parser.parse_args()

    input_path = args.input
    if not os.path.exists(input_path):
        print(f'错误：文件不存在 {input_path}')
        sys.exit(1)

    ext = Path(input_path).suffix.lower()
    if ext not in SUPPORTED_EXTS:
        print(f'错误：不支持的格式 {ext}')
        print(f'支持: {", ".join(sorted(SUPPORTED_EXTS))}')
        sys.exit(1)

    if not args.output:
        base = Path(input_path).stem
        args.output = str(Path(input_path).parent / f'{base}.md')

    ok, msg = convert_file(input_path, args.output, with_summary=not args.no_summary)

    if not args.quiet:
        print(msg)
    if ok:
        print(f'输出: {os.path.abspath(args.output)}')
    else:
        sys.exit(1)


if __name__ == '__main__':
    main()
